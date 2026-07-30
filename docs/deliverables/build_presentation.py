"""Build the SULI final presentation (.pptx).

**Slide order follows the NETL template supplied with the deliverables**
(`Fuentes Rosales, Jonathan_SULI_Presentation 2026.pdf`), because the guidelines
state that a program-provided template must be used: Title, Disclaimer,
Motivation, Background, Research Objectives, Data and Methods, Approach,
Results I, Results II, Discussion, Conclusions and Future Work, Acknowledgments.

Formatting rules from `GUIDELINES_deliverables.pdf`:

* the disclaimer appears **before any technical content** — slide 2;
* body text is **fully justified**;
* **0.5-inch left and right margins**;
* paced for a 10-minute talk plus 5 minutes of questions, with per-slide timings
  in the speaker notes.

The audience is scientists who are not specialists in load forecasting, so the
architecture is given one sentence and never unpacked.

Every number is read from `figure_facts.json` or quoted from the committed
analysis documents. The disclaimer is extracted verbatim from the supplied
template rather than retyped.

    PYTHONPATH=/tmp/pptxlib python docs/deliverables/build_presentation.py
"""

# ruff: noqa: ISC004  -- slide copy is written as adjacent string literals for
# readability; wrapping each in parentheses would obscure the prose.
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[2]
DELIV = REPO / "docs" / "deliverables"
OUT = DELIV / "Fuentes Rosales, Jonathan_SULI_Final Presentation.pptx"
FIGS = REPO / "docs" / "track_a" / "figures"
FACTS = json.loads((FIGS / "figure_facts.json").read_text())
DISCLAIMER = Path("/tmp/disclaimer_clean.txt").read_text().strip()

MARGIN = Inches(0.5)  # guidelines: 0.5-inch left and right margins
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
BODY_W = SLIDE_W - 2 * MARGIN

INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x52, 0x51, 0x4E)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, text, sub=None, size=30):
    box = slide.shapes.add_textbox(MARGIN, Inches(0.35), BODY_W, Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = INK
    r.font.name = "Calibri"
    if sub:
        q = tf.add_paragraph()
        q.text = sub
        rr = q.runs[0]
        rr.font.size = Pt(15)
        rr.font.color.rgb = MUTED
        rr.font.name = "Calibri"


def _body(slide, lines, top=1.65, size=19, width=None, left=None, space=9):
    box = slide.shapes.add_textbox(
        left or MARGIN, Inches(top), width or BODY_W, Inches(4.9)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bold = line.startswith("**")
        p.text = line.replace("**", "")
        p.alignment = PP_ALIGN.JUSTIFY  # guidelines: full justification
        p.space_after = Pt(space)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = INK
            r.font.bold = bold
            r.font.name = "Calibri"
    return box


def _note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _footer(slide, text):
    box = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.5), BODY_W, Inches(0.32))
    p = box.text_frame.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED
    r.font.name = "Calibri"


def _picture(slide, name, top=1.55, height=5.15):
    pic = slide.shapes.add_picture(
        str(FIGS / name), Emu(0), Inches(top), height=Inches(height)
    )
    pic.left = Emu(int((SLIDE_W - pic.width) / 2))
    return pic


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    cov_x = FACTS["reliability"]["extreme"]
    cov_n = FACTS["reliability"]["normal"]
    case = FACTS["case_study"]
    forest = FACTS["forest"]

    # 1 — Title
    s = _blank(prs)
    box = s.shapes.add_textbox(MARGIN, Inches(1.9), BODY_W, Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Advancing Frontiers of Energy Forecasting with Novel Methods"
    p.runs[0].font.size = Pt(40)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = INK
    p.runs[0].font.name = "Calibri"
    _body(
        s,
        [
            "Can a forecast tell you when to stop trusting it?",
            "",
            "Jonathan Fuentes Rosales",
            "Mentor: John Brewer  |  Federal Team Supervisor: Luciane Cunha",
            "Science Undergraduate Laboratory Internship (SULI), Summer 2026",
            "National Energy Technology Laboratory  |  U.S. Department of Energy",
        ],
        top=3.7,
        size=18,
        space=4,
    )
    _note(s, "0:00–0:20. Name, mentor, supervisor, program.")

    # 2 — Disclaimer, before any technical content
    s = _blank(prs)
    _title(s, "Disclaimer", size=28)
    b = _body(s, [DISCLAIMER], top=1.5, size=13, space=0)
    for p in b.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = MUTED
    _note(
        s,
        "Mandatory, and must precede any technical content. Text is verbatim "
        "from the supplied NETL template. Do not linger — read the slide "
        "title aloud and move on.",
    )

    # 3 — Motivation
    s = _blank(prs)
    _title(
        s, "Motivation", "Why accurate energy forecasting matters to the DOE mission"
    )
    _body(
        s,
        [
            "Electricity cannot be stored at grid scale, so supply and demand must be "
            "matched every second — which means someone must forecast tomorrow's "
            "demand today.",
            "",
            "**That forecast is a bet with two ways to lose.** Too low, and there is "
            "not enough power when people need it. Too high, and customers pay for "
            "generators that sat idle.",
            "",
            "Operators can hedge by committing generation early, ahead of the market. "
            "That is insurance, and it costs money — so whether it is worth buying "
            "depends entirely on how much the forecast's uncertainty can be trusted.",
            "",
            "**February 2021: forecasts missed, reserves fell short, and ERCOT cut "
            "power to paying customers.**",
        ],
    )
    _note(
        s,
        "0:20–2:00. Make it physical. Storage is the constraint; the "
        "two-sided cost is the job. Land on: trusting the uncertainty is "
        "what the hedge decision rests on. That's the thread for the talk.",
    )

    # 4 — Background
    s = _blank(prs)
    _title(
        s,
        "Background: Conventional Forecasting and Its Limits",
        "Where existing approaches fall short",
    )
    _body(
        s,
        [
            "**Current practice.** Day-ahead load forecasting relies on statistical "
            "and machine-learning models fitted to long histories of ordinary "
            "operation, assuming the relationship between weather and demand is "
            "stable.",
            "",
            "**Known limitation 1 — the assumption breaks in extremes.** During "
            "severe cold, that relationship shifts, so ordinary days actively "
            "mislead.",
            "",
            "**Known limitation 2 — the extreme days are scarce.** Twenty-four years "
            "of ERCOT data contain seventeen usable extreme cold events. A model with "
            "thousands of ordinary days has almost nothing to learn from for the days "
            "that decide whether the grid holds.",
            "",
            "**The gap this project addresses:** whether a method designed for "
            "learning from few examples helps in exactly that regime.",
        ],
        size=18,
    )
    _note(
        s,
        "2:00–3:40. The 17-events number is the constraint everything else "
        "follows from — say it slowly and let it land.",
    )

    # 5 — Research Objectives
    s = _blank(prs)
    _title(
        s, "Research Objectives", "The specific question this project set out to answer"
    )
    _body(
        s,
        [
            "**Research question.** A 2026 paper proposes AdaCNP — a model that "
            "forecasts by referring to similar past days, learning to weight past "
            "freezes more heavily than past Tuesdays. It was demonstrated on the PJM "
            "and ISO-NE grids. Does that advantage hold on Texas?",
            "",
            "**Objectives**",
            "1.  Rebuild the method and a standard baseline, and evaluate both on "
            "held-out ERCOT extreme cold events.",
            "2.  Fix the measurement, the statistical test and the success threshold "
            "in advance, committed to version control before running the experiment.",
            "3.  Establish what a study of this size can and cannot detect.",
            "",
            "**Scope.** Two models, one comparison, retrospective data. Nothing was "
            "deployed operationally.",
        ],
        size=18,
    )
    _note(
        s,
        "3:40–5:10. The pre-registration line is the most trustworthy thing "
        "you can say: 'I wrote down what would count as success before I "
        "looked at the answer.' One sentence, then move.",
    )

    # 6 — Data and Methods
    s = _blank(prs)
    _title(
        s, "Data and Methods", "What went in, what was used, and how it was evaluated"
    )
    _body(
        s,
        [
            "**Data.** Hourly ERCOT system load, 2002–2026, adopted by content hash. "
            "Seventeen load-eligible extreme cold events, defined by a regional "
            "temperature index. Inputs are calendar terms, 24 hours of recent load, "
            "and 24 hours of recent temperature — 57 features.",
            "",
            "**Guarding against cheating.** Each event is held out in turn, with a "
            "seven-day buffer on each side excluded from both training and "
            "retrieval. Only information available by 09:00 the previous day may "
            "enter a forecast. Both models receive byte-identical context sets.",
            "",
            "**Evaluation.** Probabilistic score (Gaussian negative log likelihood) "
            "against the standard-CNP baseline, paired within each event, averaged "
            "over three random seeds. 408 training runs in total.",
        ],
        size=17,
    )
    _note(
        s,
        "5:10–6:20. Don't explain neural processes architecturally. The "
        "leakage controls are worth 20 seconds — they're why the result is "
        "credible.",
    )

    # 7 — Approach
    s = _blank(prs)
    _title(s, "Approach", "End-to-end workflow from raw data to validated result")
    _body(
        s,
        [
            "1.  Verify and import the data by content hash, so provenance is exact.",
            "2.  Define episodes: predict one day's 24-hour load from information "
            "available at the previous morning's cutoff.",
            "3.  Hold out one extreme event at a time, with buffers excluded from "
            "training and from the pool of reference days.",
            "4.  **Pre-register the analysis** — endpoint, test, threshold and "
            "decision rule committed before any confirmatory run.",
            "5.  Run the full sweep: 17 events × 2 models × 3 seeds × 2 feature sets "
            "× 2 reference-selection methods = 408 runs.",
            "6.  Apply the pre-registered test exactly as written, and report what it "
            "could not have detected alongside what it did.",
        ],
        size=18,
    )
    _note(s, "6:20–7:00. Move briskly. Step 4 is the one that matters.")

    # 8 — Results I: calibration (the headline)
    s = _blank(prs)
    _title(
        s,
        "Results I: The forecasts are confidently wrong when it matters",
        "Measured on held-out extreme events, not placeholder data",
    )
    _picture(s, "fig1_reliability.png", top=1.45, height=4.75)
    _body(
        s,
        [
            f"A “90% confidence” range should contain the real load 9 times in 10. "
            f"On normal days it contains it {cov_n:.0%} of the time. During extreme "
            f"cold events, {cov_x:.0%} — and the curve sits below the line, meaning "
            f"real demand runs higher than the model expects.",
        ],
        top=6.3,
        size=15,
        space=0,
    )
    _footer(s, "Exploratory analysis")
    _note(
        s,
        "7:00–8:15. THE headline for this audience. Two numbers: 96 and 63. "
        "Then the direction — confidently wrong, and wrong the dangerous "
        "way. This is the slide to slow down on.",
    )

    # 9 — Results II: case study
    s = _blank(prs)
    _title(
        s,
        "Results II: Predicted vs. Observed Demand",
        "One day of the February 2021 event",
    )
    _picture(s, "fig2_february_2021.png", top=1.45, height=4.6)
    _body(
        s,
        [
            f"{case['day']}: all {case['hours_outside']} of {case['hours']} hours "
            f"fell outside the 90% range, with real demand roughly 10 GW above the "
            f"forecast. No load shedding occurred that day — this is forecast error, "
            f"not curtailed demand.",
        ],
        top=6.2,
        size=15,
        space=0,
    )
    _footer(s, "Exploratory · worst day of that event, selected by a stated rule")
    _note(
        s,
        "8:15–8:50. Concrete. Be honest that it's the worst day chosen by a "
        "rule — the honesty is worth more than the drama.",
    )

    # 10 — Discussion
    s = _blank(prs)
    _title(s, "Discussion", "What the results mean — and what they do not")
    _body(
        s,
        [
            f"**The method did help.** AdaCNP beat the baseline on the pre-registered "
            f"test: {forest['pooled']:+.2f} in the primary score, 95% CI "
            f"[+0.11, +0.81], p = 0.014, across all 17 events.",
            "",
            "**Three things that must be said with it.** The effect is right at the "
            "edge of what 17 events can resolve — about 77% power at the size found. "
            "It appears only in the configuration committed in advance; another "
            "configuration favoured the baseline. And it is not the effect size the "
            "original paper reported — detecting that would need several hundred "
            "events, and Texas has seventeen.",
            "",
            "**A puzzle.** The method is supposed to work by attending to similar "
            "past days. Measured directly, it does the opposite: 7.9% of its "
            "attention on cold past days, where uniform weighting would give 10.5%.",
        ],
        size=17,
    )
    _note(
        s,
        "8:50–9:40. Result, then caveats, undiluted. The caveats are why "
        "the result is believable. The mechanism puzzle is a 15-second "
        "aside — don't over-explain it.",
    )

    # 11 — Conclusions and Future Work
    s = _blank(prs)
    _title(
        s,
        "Conclusions and Future Work",
        "Key takeaways from the summer and where the work goes next",
    )
    _body(
        s,
        [
            "**Key findings.** The method helps, modestly and near the detection "
            "limit. More consequentially, both models are badly overconfident during "
            "extreme events, and predictably so — calibration degrades as the cold "
            "deepens.",
            "",
            "**Impact.** A forecast that is uncertain and says so is manageable. A "
            "forecast that is wrong and confident is what makes an operator skip the "
            "expensive hedge. Because the unreliability is predictable, a model can "
            "flag when its own range should not be trusted — directly usable for an "
            "early-commitment decision.",
            "",
            "**Future work.** Fix calibration before changing architecture; acquire "
            "day-ahead forecast weather; broaden the event definition to raise the "
            "sample size.",
        ],
        size=17,
    )
    _note(
        s,
        "9:40–10:00. Land on the 'uncertain and says so' line, then stop "
        "and take questions.",
    )

    # 12 — Acknowledgments
    s = _blank(prs)
    _title(s, "Acknowledgments", "Thank you — questions welcome")
    _body(
        s,
        [
            "**With thanks to**",
            "John Brewer — NETL, for project guidance and mentorship",
            "Luciane Cunha — Federal Team Supervisor",
            "[Add colleagues or group members who contributed]",
            "",
            "**Support.** This work was supported in part by the U.S. Department of "
            "Energy, Office of Science, Office of Workforce Development for Teachers "
            "and Scientists (WDTS), under the Science Undergraduate Laboratory "
            "Internships (SULI) program.",
            "",
            "**Questions?**  [add your email]",
        ],
        size=17,
    )
    _note(s, "Q&A slide — leave it up for the 5-minute question period.")

    # 13–14 — Backup
    s = _blank(prs)
    _title(s, "Backup: per-event results", size=26)
    _picture(s, "fig4_forest.png", top=1.2, height=5.4)
    _footer(s, "Exploratory · backup, not part of the 10-minute talk")
    _note(
        s,
        "Q&A: 'did one event drive it?' No — leaving any single event out "
        "keeps the sign, and 11 of 17 favour AdaCNP.",
    )

    s = _blank(prs)
    _title(s, "Backup: calibration vs. event severity", size=26)
    _picture(s, "fig3_coverage_vs_severity.png", top=1.2, height=5.4)
    _footer(s, "Exploratory · backup, not part of the 10-minute talk")
    _note(
        s,
        "Q&A: 'is it worse for bigger events?' Yes — about 3.5 percentage "
        "points of coverage lost per degree C of cold margin.",
    )

    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {path.name} ({path.stat().st_size:,} bytes)")
